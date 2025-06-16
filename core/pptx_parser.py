from pptx import Presentation
import os
from io import BytesIO
import uuid

from core.parser import FileParser
from core.schema import TextProperty, ImageProperty, TableProperty, BaseProperty, FileBaseProperty
from core.utils import num_tokens_from_string
from core.ocr_utils import magic_ocr
from uuid import uuid4

class PptxParser(FileParser):

    def __init__(self, file_path = None):
        super().__init__(file_path)
        self.cache_dir = f'/tmp/pptx_parser_cache/{uuid.uuid4()}'

    def extract_slide_title(self, slide):
        """提取幻灯片标题"""
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == 1:  # 标题占位符
                if shape.has_text_frame and shape.text_frame.text:
                    return shape.text_frame.text
        return ""

    def extract_text_from_shape(self, shape) -> list[BaseProperty]:
        """从形状中提取文本内容"""
        if shape.has_text_frame:
            text = shape.text_frame.text
            property_list = []
            for para in shape.text_frame.paragraphs:
                if para.text.strip():
                    text_property = TextProperty(name=f"Text from {shape.name or 'Shape'}")
                    text_property.text_content = para.text.strip()
                    text_property.page_idx = 0
                    text_property.text_level = para.level
                    text_property.content_token_length = num_tokens_from_string(text_property.text_content)
                    property_list.append(text_property)
            return property_list
        return []

    def extract_table(self, table) -> TableProperty:
        """从表格中提取数据并转换为HTML"""
        pure_text = ""
        html = "<table border='1'>"
        for row in table.rows:
            html += "<tr>"
            for cell in row.cells:
                text = cell.text_frame.text if cell.text_frame else ""
                html += f"<td>{text}</td>"
                pure_text += text + " "
            html += "</tr>"
        html += "</table>"
        table_property = TableProperty(name=f"Table {uuid.uuid4()}")
        table_property.html_content = html
        table_property.table_row_count = len(table.rows)
        table_property.table_column_count = len(table.columns)
        table_property.text_content = pure_text.strip()
        table_property.content_token_length = num_tokens_from_string(pure_text.strip())
        return table_property

    def parse(self, file_path: str) -> 'FileBaseProperty':
        presentation = Presentation(file_path)
        self.property.page_count = len(presentation.slides)
        for slide_number, slide in enumerate(presentation.slides):
            title = self.extract_slide_title(slide)
            if len(title) > 0:
                title_property = TextProperty(name=f"Slide {slide_number + 1} Title")
                title_property.text_content = title
                title_property.page_idx = slide_number
                title_property.content_token_length = num_tokens_from_string(title)
                self.property.content_list.append(title_property)
                self.property.total_text_length += len(title)
                self.property.total_token_length += title_property.content_token_length
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_properties = self.extract_text_from_shape(shape)
                    # assign page index to each shape property
                    for prop in shape_properties:
                        prop.page_idx = slide_number
                    self.property.content_list.extend(shape_properties)
                    for prop in shape_properties:
                        self.property.total_text_length += len(prop.text_content)
                        self.property.total_token_length += prop.content_token_length
                elif shape.has_table:
                    table_property = self.extract_table(shape.table)
                    table_property.page_idx = slide_number
                    self.property.content_list.append(table_property)
                    self.property.total_text_length += len(table_property.text_content)
                    self.property.total_token_length += table_property.content_token_length
                elif shape.shape_type == 13:  # 图片类型
                    img_property = ImageProperty(name=f"Image from {shape.name or 'Shape'}")
                    if shape.image:
                        img_stream = BytesIO(shape.image.blob)
                        img_property.blob_data = img_stream.getvalue()
                        img_property.image_format = shape.image.ext
                        img_property.image_path = f'{uuid.uuid4()}.{img_property.image_format}'
                        # currently, we cannot get the image width and height from pptx directly
                        
                        cache_file_path = f'{self.cache_dir}/{img_property.image_path}'
                        os.makedirs(self.cache_dir, exist_ok=True)
                        with open(cache_file_path, 'wb') as f:
                            f.write(img_property.blob_data)
                        img_property.ocr_content_list = magic_ocr(cache_file_path)
                        img_property.image_width = 0
                        img_property.image_height = 0
                        img_property.page_idx = slide_number
                        self.property.content_list.append(img_property)
                else:
                    # 其他类型的形状可以根据需要处理
                    continue
        
        return self.property